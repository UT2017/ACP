#!/usr/bin/env python
# coding: utf-8

# In[13]:


import openpyxl as xl
import pptx
import os
import platform

import tkinter as tk
import pptx.chart.data as chdata

from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt


# In[14]:


label_count = 11
slide_count_in_order_check = 1
slide_count_in_update_ppt = 1


# In[15]:


# Only these below 29 charts are supported
# Remaining 44 charts are not supported

supported_charts = [
    XL_CHART_TYPE.AREA,
    XL_CHART_TYPE.AREA_STACKED,
    XL_CHART_TYPE.AREA_STACKED_100,
    XL_CHART_TYPE.BAR_CLUSTERED,
    XL_CHART_TYPE.BAR_STACKED,
    XL_CHART_TYPE.BAR_STACKED_100,
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    XL_CHART_TYPE.COLUMN_STACKED,
    XL_CHART_TYPE.COLUMN_STACKED_100,
    XL_CHART_TYPE.BUBBLE,
    XL_CHART_TYPE.BUBBLE_THREE_D_EFFECT,
    XL_CHART_TYPE.DOUGHNUT,
    XL_CHART_TYPE.DOUGHNUT_EXPLODED,
    XL_CHART_TYPE.LINE,
    XL_CHART_TYPE.LINE_MARKERS,
    XL_CHART_TYPE.LINE_MARKERS_STACKED,
    XL_CHART_TYPE.LINE_MARKERS_STACKED_100,
    XL_CHART_TYPE.LINE_STACKED,
    XL_CHART_TYPE.LINE_STACKED_100,
    XL_CHART_TYPE.PIE,
    XL_CHART_TYPE.PIE_EXPLODED,
    XL_CHART_TYPE.RADAR,
    XL_CHART_TYPE.RADAR_FILLED,
    XL_CHART_TYPE.RADAR_MARKERS,
    XL_CHART_TYPE.XY_SCATTER,
    XL_CHART_TYPE.XY_SCATTER_LINES,
    XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
    XL_CHART_TYPE.XY_SCATTER_SMOOTH,
    XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS,
]


# In[16]:


def check_order_of_multiple_charts_helper():
    global label_count

    try:
        check_order_of_multiple_charts()

    except Exception as ex:
        template = "An exception of type {0} occurred. Arguments:\n{1!r}"
        message = template.format(type(ex).__name__, ex.args)
        message1 = "In slide no. " + str(slide_count_in_order_check)
        tk.Label(
            frame,
            anchor=tk.W,
            justify=tk.LEFT,
            font="Helvetica 20 bold",
            foreground="red",
            bg="#ffffff",
            text="\n" + message1,
        ).grid(row=label_count, sticky=tk.W, padx=(4, 4))
        label_count = label_count + 1
        tk.Label(
            frame,
            anchor=tk.W,
            justify=tk.LEFT,
            font="Helvetica 20 bold",
            foreground="red",
            bg="#ffffff",
            text=message + "\n",
        ).grid(row=label_count, sticky=tk.W, padx=(4, 4))
        label_count = label_count + 1


def update_ppt_helper():
    global label_count
    var1 = 0

    try:
        update_ppt()
    except Exception as ex:
        template = "An exception of type {0} occurred. Arguments:\n{1!r}"
        message = template.format(type(ex).__name__, ex.args)
        message1 = "In slide no. " + str(slide_count_in_update_ppt)
        tk.Label(
            frame,
            anchor=tk.W,
            justify=tk.LEFT,
            font="Helvetica 20 bold",
            foreground="red",
            bg="#ffffff",
            text="\n" + message1,
        ).grid(row=label_count, sticky=tk.W, padx=(4, 4))
        label_count = label_count + 1
        tk.Label(
            frame,
            anchor=tk.W,
            justify=tk.LEFT,
            font="Helvetica 20 bold",
            foreground="red",
            bg="#ffffff",
            text=message + "\n",
        ).grid(row=label_count, sticky=tk.W, padx=(4, 4))
        label_count = label_count + 1
        var1 = 1
    finally:
        if var1 == 0:
            tk.Label(
                frame,
                anchor=tk.W,
                justify=tk.LEFT,
                font="Helvetica 20 bold",
                foreground="green",
                bg="#ffffff",
                text="Success!\nPlease find updated ppt at " + os.getcwd(),
            ).grid(row=label_count, sticky=tk.W, padx=(4, 4))
            label_count = label_count + 1


# In[ ]:


def check_order_of_multiple_charts():
    global label_count
    global slide_count_in_order_check

    input_master_excel_name_xlsx = input_master_excel_name.get() + ".xlsx"
    input_ppt_name_pptx = input_ppt_name.get() + ".pptx"

    try:
        presentation_object = pptx.Presentation(input_ppt_name_pptx)
        slides_object = presentation_object.slides

    except:
        tk.Label(
            frame,
            anchor=tk.W,
            justify=tk.LEFT,
            font="Helvetica 20 bold",
            foreground="red",
            bg="#ffffff",
            text="Input ppt not found.\nMake sure to give exact name without the extension",
        ).grid(row=label_count, sticky=tk.W, padx=(4, 4))
        label_count = label_count + 1

    for slides_iterator in slides_object:

        for shapes_iterator in slides_iterator.shapes:

            if shapes_iterator.has_chart:
                graphic_frame = shapes_iterator
                chart = graphic_frame.chart

                try:

                    if shapes_iterator.chart.chart_type in supported_charts:

                        tk.Label(
                            frame,
                            anchor=tk.W,
                            justify=tk.LEFT,
                            font="Helvetica 12 bold",
                            bg="#ffffff",
                            text="Slide no. "
                            + str(slide_count_in_order_check)
                            + "\n    Chart type "
                            + str(shapes_iterator.chart.chart_type),
                        ).grid(row=label_count, sticky=tk.W, padx=(4, 4))
                        label_count = label_count + 1

                        chart_title = shapes_iterator.chart.chart_title
                        if chart_title.has_text_frame:
                            tk.Label(
                                frame,
                                anchor=tk.W,
                                justify=tk.LEFT,
                                font="Helvetica 12 bold",
                                bg="#ffffff",
                                text="    Chart title "
                                + chart_title.text_frame.text
                                + "\n\n",
                            ).grid(row=label_count, sticky=tk.W, padx=(4, 4))
                            label_count = label_count + 1

                        for plot in shapes_iterator.chart.plots:
                            plot.has_data_labels = True

                except Exception as ex:

                    template = "An exception of type {0} occurred. Arguments:\n{1!r}"
                    message = template.format(type(ex).__name__, ex.args)

                    tk.Label(
                        frame,
                        anchor=tk.W,
                        justify=tk.LEFT,
                        foreground="red",
                        font="Helvetica 20 bold",
                        bg="#ffffff",
                        text="Slide no. "
                        + str(slide_count_in_order_check)
                        + "\n"
                        + "An unsupported chart exists. Please refer to the unsupported charts list. You would have to manually input data for this chart.\n Please find exception details below: \n",
                    ).grid(row=label_count, sticky=tk.W, padx=(4, 4))
                    label_count = label_count + 1

                    tk.Label(
                        frame,
                        anchor=tk.W,
                        justify=tk.LEFT,
                        font="Helvetica 20 bold",
                        foreground="red",
                        bg="#ffffff",
                        text=message + "\n",
                    ).grid(row=label_count, sticky=tk.W, padx=(4, 4))
                    label_count = label_count + 1

                    continue

        slide_count_in_order_check = slide_count_in_order_check + 1

    slide_count_in_order_check = 1


# In[ ]:


def update_ppt():
    global slide_count_in_update_ppt
    global label_count

    input_master_excel_name_xlsx = input_master_excel_name.get() + ".xlsx"
    input_ppt_name_pptx = input_ppt_name.get() + ".pptx"
    output_ppt_name_pptx = output_ppt_name.get() + ".pptx"

    try:
        presentation_object = pptx.Presentation(input_ppt_name_pptx)
    except:
        tk.Label(
            frame,
            anchor=tk.W,
            justify=tk.LEFT,
            font="Helvetica 20 bold",
            foreground="red",
            bg="#ffffff",
            text="Input ppt not found.\nMake sure to give exact name without the extension",
        ).grid(row=label_count, sticky=tk.W, padx=(4, 4))
        label_count = label_count + 1

    try:
        workbook_object = xl.load_workbook(input_master_excel_name_xlsx)
    except:
        tk.Label(
            frame,
            anchor=tk.W,
            justify=tk.LEFT,
            font="Helvetica 20 bold",
            foreground="red",
            bg="#ffffff",
            text="Input master excel not found.\nMake sure to give exact name without the extension",
        ).grid(row=label_count, sticky=tk.W, padx=(4, 4))
        label_count = label_count + 1

    slides_object = presentation_object.slides

    for slides_iterator in slides_object:
        charts_count = 1
        for shapes_iterator in slides_iterator.shapes:
            if shapes_iterator.has_chart:
                sheet_object = workbook_object.get_sheet_by_name(
                    str(slide_count_in_update_ppt) + "_" + str(charts_count)
                )
                charts_count = charts_count + 1

                try:

                    # Handling Scatter plots
                    if (
                        shapes_iterator.chart.chart_type == XL_CHART_TYPE.XY_SCATTER
                        or shapes_iterator.chart.chart_type
                        == XL_CHART_TYPE.XY_SCATTER_LINES
                        or shapes_iterator.chart.chart_type
                        == XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS
                        or shapes_iterator.chart.chart_type
                        == XL_CHART_TYPE.XY_SCATTER_SMOOTH
                        or shapes_iterator.chart.chart_type
                        == XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS
                    ):

                        # Different chart object type to be created
                        chart_data_object = chdata.XyChartData()

                        rows_in_current_sheet = list(sheet_object.rows)
                        rows_in_current_sheet.pop(0)

                        series_object = chart_data_object.add_series("series 0")

                        for rows_iterator in rows_in_current_sheet:
                            rows_iterator = list(
                                rows_iterator
                            )  # converting tuple to list for easy manipulation
                            var1 = 0
                            for current_row_iterator in rows_iterator:
                                rows_iterator[var1] = current_row_iterator.value
                                var1 = var1 + 1
                            series_object.add_data_point(
                                rows_iterator[0], rows_iterator[1]
                            )

                    # Handling Bubble plots
                    elif (
                        shapes_iterator.chart.chart_type == XL_CHART_TYPE.BUBBLE
                        or shapes_iterator.chart.chart_type
                        == XL_CHART_TYPE.BUBBLE_THREE_D_EFFECT
                    ):

                        # Different chart object type to be created
                        chart_data_object = chdata.BubbleChartData()

                        rows_in_current_sheet = list(sheet_object.rows)
                        rows_in_current_sheet.pop(0)

                        series_object = chart_data_object.add_series("series 0")

                        for rows_iterator in rows_in_current_sheet:
                            rows_iterator = list(
                                rows_iterator
                            )  # converting tuple to list for easy manipulation
                            var1 = 0
                            for current_row_iterator in rows_iterator:
                                rows_iterator[var1] = current_row_iterator.value
                                var1 = var1 + 1
                            series_object.add_data_point(
                                rows_iterator[0], rows_iterator[1], rows_iterator[2]
                            )

                    elif shapes_iterator.chart.chart_type in supported_charts:
                        # Different chart object type to be created
                        chart_data_object = chdata.CategoryChartData()

                        columns_in_current_sheet = list(sheet_object.columns)

                        first_column_in_current_sheet = list(
                            columns_in_current_sheet[0]
                        )
                        first_column_in_current_sheet.pop(0)

                        # Adding categories using the first column
                        for (
                            first_column_in_current_sheet_iterator
                        ) in first_column_in_current_sheet:
                            chart_data_object.add_category(
                                first_column_in_current_sheet_iterator.value
                            )

                        columns_in_current_sheet.pop(0)

                        for columns_iterator in columns_in_current_sheet:
                            columns_iterator = list(
                                columns_iterator
                            )  # converting tuple to list for easy manipulation
                            series_name = columns_iterator.pop(0)
                            var1 = 0
                            for current_column_iterator in columns_iterator:
                                columns_iterator[var1] = current_column_iterator.value
                                var1 = var1 + 1

                            # converting back to tuple since add_series function takes tuples as values
                            columns_iterator = tuple(columns_iterator)

                            chart_data_object.add_series(
                                series_name.value, columns_iterator
                            )

                    shapes_iterator.chart.replace_data(chart_data_object)

                    for plot in shapes_iterator.chart.plots:
                        plot.has_data_labels = True

                except Exception as ex:

                    template = "An exception of type {0} occurred. Arguments:\n{1!r}"
                    message = template.format(type(ex).__name__, ex.args)

                    tk.Label(
                        frame,
                        anchor=tk.W,
                        justify=tk.LEFT,
                        foreground="red",
                        font="Helvetica 20 bold",
                        bg="#ffffff",
                        text="In Slide no. "
                        + str(slide_count_in_update_ppt)
                        + "\n"
                        + "An unsupported chart exists. Please refer to the unsupported charts list.\n You would have to manually input data for this chart. Skipping this chart to go populate the next.\n Please find exception details below: \n",
                    ).grid(row=label_count, sticky=tk.W, padx=(4, 4))
                    label_count = label_count + 1

                    tk.Label(
                        frame,
                        anchor=tk.W,
                        justify=tk.LEFT,
                        font="Helvetica 20 bold",
                        foreground="red",
                        bg="#ffffff",
                        text=message + "\n",
                    ).grid(row=label_count, sticky=tk.W, padx=(4, 4))
                    label_count = label_count + 1

                    continue

        slide_count_in_update_ppt = slide_count_in_update_ppt + 1

    presentation_object.save(output_ppt_name_pptx)
    slide_count_in_update_ppt = 1


# In[ ]:


"""GUI"""


# In[ ]:


main_window = tk.Tk()
main_window.configure(background="#ffffff")
main_window.title("Automating Charts population")


# In[ ]:


def onFrameConfigure(canvas):
    """Reset the scroll region to encompass the inner frame"""
    canvas.configure(scrollregion=canvas.bbox("all"))


# In[ ]:


def _on_mousewheel(event):
    canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")


# In[ ]:


canvas = tk.Canvas(main_window, borderwidth=0, bg="#ffffff")
frame = tk.Frame(canvas, bg="#ffffff")

# scrollbar
vsb = tk.Scrollbar(main_window, orient="vertical", command=canvas.yview)
canvas.configure(yscrollcommand=vsb.set)
hsb = tk.Scrollbar(main_window, orient="horizontal", command=canvas.xview)
canvas.configure(xscrollcommand=hsb.set)
canvas.bind_all("<MouseWheel>", _on_mousewheel)

hsb.pack(side="bottom", fill="x")
vsb.pack(side="right", fill="y")
canvas.pack(side="left", fill="both", expand=True)
canvas.create_window((0, 0), window=frame, anchor="nw")

frame.bind("<Configure>", lambda event, canvas=canvas: onFrameConfigure(canvas))


# In[ ]:


tk.Label(
    frame,
    anchor=tk.W,
    justify=tk.LEFT,
    font="Helvetica 20 bold",
    foreground="#ffffff",
    bg="#ff7919",
    text="Populating powerpoint charts made easy",
).grid(sticky=(tk.N, tk.S, tk.E, tk.W), rowspan=2, columnspan=3)

# Printing 1 empty row
tk.Label(frame, anchor=tk.W, justify=tk.LEFT, bg="#ffffff").grid(
    row=2, sticky=(tk.N, tk.S, tk.E, tk.W)
)

tk.Label(
    frame,
    anchor=tk.W,
    justify=tk.LEFT,
    font="Helvetica 12 ",
    bg="#ffffff",
    text="Input Master excel file name without the extension(.xlsx)",
).grid(row=3, sticky=(tk.N, tk.S, tk.E, tk.W), padx=(4, 4))
tk.Label(
    frame,
    anchor=tk.W,
    justify=tk.LEFT,
    font="Helvetica 12 ",
    bg="#ffffff",
    text="Input powerpoint file name without the extension(.pptx)",
).grid(row=4, sticky=(tk.N, tk.S, tk.E, tk.W), padx=(4, 4))
tk.Label(
    frame,
    anchor=tk.W,
    justify=tk.LEFT,
    font="Helvetica 12 ",
    bg="#ffffff",
    text="Input updated (output) powerpoint file name without the extension(.pptx)",
).grid(row=5, sticky=(tk.N, tk.S, tk.E, tk.W), padx=(4, 4))

# Printing 3 empty rows
tk.Label(frame, anchor=tk.W, justify=tk.LEFT, bg="#ffffff").grid(
    row=6, sticky=(tk.N, tk.S, tk.E, tk.W)
)
tk.Label(frame, anchor=tk.W, justify=tk.LEFT, bg="#ffffff").grid(
    row=8, sticky=(tk.N, tk.S, tk.E, tk.W)
)
tk.Label(frame, anchor=tk.W, justify=tk.LEFT, bg="#ffffff").grid(
    row=9, sticky=(tk.N, tk.S, tk.E, tk.W)
)

frame.columnconfigure(0, weight=100, minsize=700)
frame.columnconfigure(1, weight=50, minsize=600)
frame.columnconfigure(2, weight=3000, minsize=600)

input_master_excel_name = tk.Entry(
    frame, selectborderwidth=100, relief="sunken", width=60
)
input_ppt_name = tk.Entry(frame, selectborderwidth=100, relief="sunken", width=60)
output_ppt_name = tk.Entry(frame, selectborderwidth=100, relief="sunken", width=60)

input_master_excel_name.grid(row=3, column=1, sticky=tk.W)
input_ppt_name.grid(row=4, column=1, sticky=tk.W)
output_ppt_name.grid(row=5, column=1, sticky=tk.W)


if platform.system() == "Darwin":  # if its a Mac
    tk.Button(
        frame,
        font="Helvetica 12",
        width=30,
        highlightbackground="#3E4149",
        bg="#3B5998",
        fg="#ffffff",
        text="Check order of multiple charts",
        command=check_order_of_multiple_charts_helper,
        justify=tk.LEFT,
    ).grid(row=7, column=0, sticky=tk.W, padx=(4, 4))
    tk.Button(
        frame,
        font="Helvetica 12",
        width=30,
        highlightbackground="#3E4149",
        bg="#3B5998",
        fg="#ffffff",
        text="Get updated ppt",
        command=update_ppt_helper,
        justify=tk.LEFT,
    ).grid(row=7, column=1, sticky=tk.W, padx=(4, 4))
else:  # if its Windows or Linux
    tk.Button(
        frame,
        font="Helvetica 12",
        width=30,
        bg="#3B5998",
        fg="#ffffff",
        text="Check order of multiple charts",
        command=check_order_of_multiple_charts_helper,
        justify=tk.LEFT,
    ).grid(row=7, column=0, sticky=tk.W, padx=(4, 4))
    tk.Button(
        frame,
        font="Helvetica 12",
        width=30,
        bg="#3B5998",
        fg="#ffffff",
        text="Get updated ppt",
        command=update_ppt_helper,
        justify=tk.LEFT,
    ).grid(row=7, column=1, sticky=tk.W, padx=(4, 4))
    tk.Button(
        frame,
        font="Helvetica 12",
        width=30,
        bg="#3B5998",
        fg="#ffffff",
        text="Quit",
        command=main_window.destroy,
        justify=tk.LEFT,
    ).grid(row=7, column=2, sticky=tk.W, padx=(4, 4))


main_window.mainloop()
